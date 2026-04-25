"""
retrieve.py — Three-signal hybrid retrieval system
Signals: (1) Semantic embeddings via sentence-transformers
         (2) BM25 exact keyword/token matching
         (3) Glossary source weighting for variable lookups

Handles both broad queries ("main finding of study") and narrow queries
("what does frm_brand_awareness measure" or just "brand awareness").

No external LLM APIs. Runs fully locally.
Author: Yuvanesh Raju
"""

import os
import re
import json
import hashlib
import email
from pathlib import Path
from typing import List, Dict


# ─────────────────────────────────────────────────────────────────
# 1. FILE PARSERS
# ─────────────────────────────────────────────────────────────────

def parse_pdf(path: str) -> List[Dict]:
    from pypdf import PdfReader
    reader = PdfReader(path)
    chunks = []
    for page_num, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            for chunk in _chunk_text(text, max_words=120, overlap=30):
                chunks.append({
                    "text": chunk,
                    "source": os.path.basename(path),
                    "source_type": "pdf",
                    "meta": f"page {page_num}",           # FIX: page number preserved
                })
    return chunks


def parse_pptx(path: str) -> List[Dict]:
    from pptx import Presentation
    prs = Presentation(path)
    chunks = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        combined = " | ".join(parts)
        if combined.strip():
            for chunk in _chunk_text(combined, max_words=120, overlap=30):
                chunks.append({
                    "text": chunk,
                    "source": os.path.basename(path),
                    "source_type": "pptx",
                    "meta": f"slide {slide_num}",         # FIX: slide number preserved
                })
    return chunks


def parse_docx(path: str) -> List[Dict]:
    from docx import Document
    doc = Document(path)
    paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    combined = " ".join(paras)
    chunks = []
    for chunk in _chunk_text(combined, max_words=120, overlap=30):
        chunks.append({
            "text": chunk,
            "source": os.path.basename(path),
            "source_type": "docx",
            "meta": "",
        })
    return chunks


def parse_eml(path: str) -> List[Dict]:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        msg = email.message_from_file(f)
    subject = msg.get("Subject", "")
    sender  = msg.get("From", "")
    date    = msg.get("Date", "")       # FIX: date now captured
    body = ""
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                body += part.get_payload(decode=True).decode("utf-8", errors="ignore")
    else:
        payload = msg.get_payload(decode=True)
        if payload:
            body = payload.decode("utf-8", errors="ignore")
    full_text = f"Subject: {subject}\nFrom: {sender}\nDate: {date}\n\n{body}"
    chunks = []
    for chunk in _chunk_text(full_text, max_words=120, overlap=30):
        chunks.append({
            "text": chunk,
            "source": os.path.basename(path),
            "source_type": "eml",
            "meta": f"from: {sender} | {date}",          # FIX: structured sender/date
        })
    return chunks


def parse_json_glossary(path: str) -> List[Dict]:
    """
    Parse glossary JSON.
    Each variable gets its own chunk tagged source_type='glossary'.
    Glossary chunks receive a score boost during retrieval for variable queries.
    """
    with open(path, "r", encoding="utf-8") as f:
        data = json.load(f)
    chunks = []

    # Sheet-level definitions
    for sheet_name, info in data.get("SHEET_DEFINITIONS", {}).items():
        definition = info.get("definition", "")
        use_for    = ", ".join(info.get("use_for", []))
        text = f"Sheet '{sheet_name}': {definition} Use for: {use_for}"
        chunks.append({
            "text": text.strip(),
            "source": "glossary",
            "source_type": "glossary",
            "meta": f"sheet: {sheet_name}",
        })

    # Variable definitions — one chunk per variable (already atomic)
    for var_name, info in data.get("V", {}).items():
        if var_name == "attribute_type":
            continue
        if isinstance(info, dict):
            definition = info.get("definition", "")
            var_type   = info.get("type", "")
            table      = info.get("table", "")
            text = (
                f"Variable '{var_name}': {definition}"
                + (f" Type: {var_type}." if var_type else "")
                + (f" Table: {str(table)}." if str(table) != "" else "")
            )
        else:
            text = f"Variable '{var_name}': {info}"
        chunks.append({
            "text": text.strip(),
            "source": "glossary",
            "source_type": "glossary",
            "meta": f"variable: {var_name}",
        })

    return chunks


# ─────────────────────────────────────────────────────────────────
# 2. CHUNKING UTILITY
# ─────────────────────────────────────────────────────────────────

def _chunk_text(text: str, max_words: int = 120, overlap: int = 30) -> List[str]:
    """Overlapping word-window chunking."""
    words = text.split()
    if not words:
        return []
    chunks, start = [], 0
    while start < len(words):
        end   = min(start + max_words, len(words))
        chunk = " ".join(words[start:end]).strip()
        if chunk:
            chunks.append(chunk)
        if end == len(words):
            break
        start += max_words - overlap
    return chunks


# ─────────────────────────────────────────────────────────────────
# 3. LOAD ALL DOCUMENTS
# ─────────────────────────────────────────────────────────────────

def load_all_chunks(directory: str) -> List[Dict]:
    """Scan directory, dispatch to correct parser by extension."""
    dir_path = Path(directory)
    parsers  = {
        ".pdf":  parse_pdf,
        ".pptx": parse_pptx,
        ".docx": parse_docx,
        ".eml":  parse_eml,
        ".json": parse_json_glossary,
    }
    chunks = []
    for file_path in sorted(dir_path.iterdir()):
        ext = file_path.suffix.lower()
        if ext in parsers:
            try:
                chunks.extend(parsers[ext](str(file_path)))
            except Exception as e:
                print(f"[WARN] Could not parse {file_path.name}: {e}")
    return chunks


# ─────────────────────────────────────────────────────────────────
# 4. QUERY CLASSIFIER
# ─────────────────────────────────────────────────────────────────

# Variable-name pattern: 2+ underscore-joined segments of lowercase alphanumeric
_VAR_PATTERN = re.compile(r"\b[a-z][a-z0-9]*(?:_[a-z0-9]+){1,}\b")

_NARROW_KEYWORDS = re.compile(
    r"\b(variable|variables|define|definition|what is|what does|"
    r"mean(ing)?|measure|captures?|represent|codebook|column|field)\b",
    re.IGNORECASE,
)

def _query_is_narrow(query: str) -> bool:
    """
    Returns True for variable-lookup / definition queries.
    Handles both long ('what does frm_brand_awareness measure') and
    SHORT queries ('brand awareness', 'frm_brand_awareness').
    Short queries containing an underscore variable name are always narrow.
    """
    q = query.strip()
    # Explicit variable name in query (underscored token)
    if _VAR_PATTERN.search(q.lower()):
        return True
    # Explicit definitional keywords
    if _NARROW_KEYWORDS.search(q):
        return True
    # Short query (≤4 words) — likely a lookup, not a broad research question
    if len(q.split()) <= 4:
        return True
    return False


# ─────────────────────────────────────────────────────────────────
# 5. EMBEDDING MODEL  (semantic signal)
# ─────────────────────────────────────────────────────────────────

_embed_model  = None
_embed_failed = False   # set True if model unavailable; fall back gracefully


def _get_embed_model():
    """
    Load sentence-transformers model once and cache.
    Uses all-MiniLM-L6-v2 (384-dim, 80 MB).
    Falls back gracefully if model cannot be loaded (e.g. no internet on first run).
    """
    global _embed_model, _embed_failed
    if _embed_model is not None or _embed_failed:
        return _embed_model
    try:
        from sentence_transformers import SentenceTransformer
        _embed_model = SentenceTransformer("all-MiniLM-L6-v2")
        return _embed_model
    except Exception as e:
        print(f"[INFO] Embedding model unavailable ({e}). Using lexical fallback.")
        _embed_failed = True
        return None


def _embed(texts: List[str]):
    """Return L2-normalised embeddings, or None if model unavailable."""
    import numpy as np
    model = _get_embed_model()
    if model is None:
        return None
    embs  = model.encode(texts, batch_size=64, show_progress_bar=False,
                          convert_to_numpy=True)
    norms = np.linalg.norm(embs, axis=1, keepdims=True)
    norms[norms == 0] = 1.0
    return embs / norms


# ─────────────────────────────────────────────────────────────────
# 6. BM25 SCORER  (lexical signal)
# ─────────────────────────────────────────────────────────────────

def _tokenize(text: str) -> List[str]:
    """Keep underscored tokens intact (critical for variable names)."""
    return re.findall(r"[a-z0-9_]+", text.lower())


def _bm25_scores(query: str, bm25_index) -> "np.ndarray":
    """
    Return normalised BM25 scores using a pre-built BM25Okapi index.

    FIX: BM25 index is no longer rebuilt on every query call.
    It is built once in _build_index() and passed in here.
    This avoids O(corpus) construction cost per query.
    """
    import numpy as np
    raw = np.array(bm25_index.get_scores(_tokenize(query)), dtype=np.float32)
    mx  = raw.max()
    return raw / mx if mx > 0 else raw


# ─────────────────────────────────────────────────────────────────
# 7. INDEX CACHE
# ─────────────────────────────────────────────────────────────────

_index_cache: Dict = {}


def _build_index(directory: str) -> Dict:
    """
    Parse all files, build embeddings and BM25 index once per session.

    FIX: BM25Okapi is now built here and cached alongside embeddings,
    so it is constructed exactly once per directory, not once per query.

    FIX: Glossary chunk positions are pre-computed as a boolean mask
    so the glossary boost loop in retrieve() is O(1) per chunk instead
    of doing a dict lookup inside a Python for-loop.
    """
    import numpy as np
    from rank_bm25 import BM25Okapi

    chunks = load_all_chunks(directory)
    texts  = [c["text"] for c in chunks]

    print(f"[INFO] Loaded {len(chunks)} chunks from {directory}")

    # Semantic embeddings (may be None if model unavailable)
    embs = _embed(texts)
    if embs is not None:
        print(f"[INFO] Embeddings built: shape {embs.shape}")
    else:
        print("[INFO] Running in lexical-only mode (no embeddings).")

    # BM25 index — built ONCE here, reused on every query call
    tokenized   = [_tokenize(t) for t in texts]
    bm25_index  = BM25Okapi(tokenized)
    print(f"[INFO] BM25 index built over {len(texts)} passages.")

    # Pre-compute glossary mask for O(1) boost application
    glossary_mask = np.array(
        [c["source_type"] == "glossary" for c in chunks], dtype=bool
    )

    return {
        "chunks":        chunks,
        "texts":         texts,
        "embeddings":    embs,
        "bm25":          bm25_index,       # FIX: cached BM25 index
        "glossary_mask": glossary_mask,    # FIX: pre-computed boolean mask
    }


def _get_index(directory: str) -> Dict:
    if directory not in _index_cache:
        _index_cache[directory] = _build_index(directory)
    return _index_cache[directory]


# ─────────────────────────────────────────────────────────────────
# 8. MAIN RETRIEVE FUNCTION
# ─────────────────────────────────────────────────────────────────

# Glossary source boost: added to final score for glossary chunks on narrow queries.
# Value of 0.15 was empirically tuned: large enough to surface glossary definitions
# above prose that happens to contain the same keywords, small enough not to dominate
# on queries where a document chunk is genuinely more relevant.
GLOSSARY_BOOST = 0.15


def retrieve(query: str, directory: str = "test_files") -> List[Dict]:
    """
    Return top-5 most relevant passages for the given query.

    Each result dict:
      - "text":   str   — passage content
      - "source": str   — source filename or "glossary"
      - "score":  float — composite relevance score (higher = better)
      - "meta":   str   — page/slide/sender info for display

    Three-signal hybrid:
      semantic_score  — cosine similarity of sentence embeddings (handles paraphrase)
      bm25_score      — exact token match (handles variable names, proper nouns)
      glossary_boost  — adds GLOSSARY_BOOST to glossary chunks on narrow queries

    Blend weights:
      narrow query → α=0.45 semantic + 0.55 BM25  (lean on exact match)
      broad  query → α=0.75 semantic + 0.25 BM25  (lean on meaning)

    When embeddings are unavailable (no internet on first run), falls back to
    BM25-only with the same glossary boost.
    """
    import numpy as np

    idx    = _get_index(directory)
    chunks = idx["chunks"]
    embs   = idx["embeddings"]

    if not chunks:
        return []

    narrow = _query_is_narrow(query)

    # ── Semantic signal ──────────────────────────────────────────
    if embs is not None:
        q_emb = _embed([query])
        if q_emb is not None:
            sem = embs.dot(q_emb[0])                          # cosine sim (already L2-normalised)
            sem = (sem - sem.min()) / (sem.max() - sem.min() + 1e-9)  # scale to [0,1]
        else:
            sem = np.zeros(len(chunks), dtype=np.float32)
    else:
        sem = np.zeros(len(chunks), dtype=np.float32)

    # ── BM25 signal ──────────────────────────────────────────────
    # FIX: uses pre-built cached index instead of rebuilding on every call
    bm25 = _bm25_scores(query, idx["bm25"])

    # ── Blend ────────────────────────────────────────────────────
    alpha = 0.45 if narrow else 0.75                          # weight for semantic
    if embs is None:
        alpha = 0.0                                            # fallback: BM25 only

    final = alpha * sem + (1.0 - alpha) * bm25

    # ── Glossary boost ───────────────────────────────────────────
    # On narrow queries, glossary chunks get a flat bonus so variable definitions
    # surface above generic text that happens to contain the same words.
    # FIX: uses pre-computed boolean mask — no per-chunk dict lookup needed
    if narrow:
        final = final.copy()
        final[idx["glossary_mask"]] = np.minimum(
            1.0, final[idx["glossary_mask"]] + GLOSSARY_BOOST
        )

    # ── Top-5 with deduplication ─────────────────────────────────
    # FIX: use full-text hash for dedup instead of first-80-chars prefix
    # (prefix can collide on chunks with similar openings)
    top_idx = np.argsort(final)[::-1]
    results, seen = [], set()
    for i in top_idx:
        text_hash = hashlib.md5(chunks[i]["text"].encode()).hexdigest()
        if text_hash in seen:
            continue
        seen.add(text_hash)
        results.append({
            "text":   chunks[i]["text"],
            "source": chunks[i]["source"],
            "score":  float(final[i]),
            "meta":   chunks[i].get("meta", ""),
        })
        if len(results) == 5:
            break

    return results


# ─────────────────────────────────────────────────────────────────
# 9. EVALUATION HARNESS
# ─────────────────────────────────────────────────────────────────

# Ground-truth relevance labels: (query, expected_source_keywords)
# A result is relevant if its source contains any expected keyword.
EVAL_SET = [
    # Broad queries
    ("What was the main finding of the brand study?",
     ["glossary", "research_proposal", "phase2"]),
    ("What are the cooking habits of probable trialists?",
     ["glossary", "email_01", "email_02"]),
    ("How many respondents were surveyed?",
     ["glossary", "email_01", "email_05"]),
    ("What competitors appear in the frozen meal market?",
     ["email_03", "email_04"]),
    ("What channels do consumers use to buy frozen ready meals?",
     ["email_02", "glossary"]),
    ("What was discussed in the project kickoff?",
     ["email_01"]),
    ("What is the share of wallet for regular buyers?",
     ["email_04", "glossary", "email_01"]),
    # Narrow queries
    ("What does the variable frm_brand_awareness measure?",
     ["glossary"]),
    ("What does N_Home_Cooking_per_week mean?",
     ["email_02", "glossary"]),
    ("What does frm_brand_consideration capture?",
     ["glossary"]),
    ("What are probable_trialists?",
     ["email_01", "glossary"]),
    ("What is total_count?",
     ["glossary"]),
    ("What does marketing_dimension represent?",
     ["email_03", "email_04"]),
    ("brand awareness",                                   # short narrow query
     ["glossary", "email_02"]),
]


def _is_relevant(result: Dict, expected_sources: List[str]) -> bool:
    src = result["source"].lower()
    return any(kw in src for kw in expected_sources)


def evaluate(directory: str = "test_files") -> Dict:
    """
    Compute MRR and P@1 over the eval set.
    Returns dict with per-query detail and aggregate metrics.
    """
    mrr_sum, p1_sum = 0.0, 0.0
    rows = []

    for query, expected in EVAL_SET:
        results = retrieve(query, directory)
        # MRR: reciprocal rank of first relevant result in top-5
        rr = 0.0
        for rank, r in enumerate(results, 1):
            if _is_relevant(r, expected):
                rr = 1.0 / rank
                break
        # P@1: is the top result relevant?
        p1 = 1.0 if results and _is_relevant(results[0], expected) else 0.0
        mrr_sum += rr
        p1_sum  += p1
        rows.append({
            "query":    query,
            "type":     "narrow" if _query_is_narrow(query) else "broad",
            "rr":       round(rr, 3),
            "p@1":      int(p1),
            "top_src":  results[0]["source"] if results else "—",
            "expected": expected,
        })

    n = len(EVAL_SET)
    return {
        "MRR":    round(mrr_sum / n, 3),
        "P@1":    round(p1_sum  / n, 3),
        "detail": rows,
    }


# ─────────────────────────────────────────────────────────────────
# 10. CLI  —  python retrieve.py
# ─────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    DEMO_QUERIES = [
        # Broad
        "What was the main finding of the brand study?",
        "What are the cooking habits of probable trialists?",
        "What channels do consumers use to buy frozen ready meals?",
        # Narrow — long
        "What does the variable frm_brand_awareness measure?",
        "What does N_Home_Cooking_per_week mean?",
        # Narrow — SHORT  (the hard case)
        "brand awareness",
        "frm_brand_consideration",
    ]

    DIR = "test_files"
    print("=" * 65)
    print("  RETRIEVAL DEMO")
    print("=" * 65)
    for q in DEMO_QUERIES:
        print(f"\nQ [{('narrow' if _query_is_narrow(q) else 'broad'):6}]: {q}")
        for i, r in enumerate(retrieve(q, DIR), 1):
            meta = f" [{r['meta']}]" if r.get("meta") else ""
            print(f"  {i}. [{r['source']:35s}]{meta} {r['score']:.4f}  {r['text'][:90]}…")

    print("\n" + "=" * 65)
    print("  EVALUATION METRICS")
    print("=" * 65)
    metrics = evaluate(DIR)
    print(f"\n  MRR  = {metrics['MRR']}")
    print(f"  P@1  = {metrics['P@1']}")
    print(f"\n  {'Query':<47} {'Type':6} {'RR':>5}  {'P@1':>3}  Top source")
    print("  " + "-" * 90)
    for row in metrics["detail"]:
        q_short = row["query"][:45] + ("…" if len(row["query"]) > 45 else "")
        print(f"  {q_short:<47} {row['type']:6} {row['rr']:>5}  {row['p@1']:>3}  {row['top_src']}")
